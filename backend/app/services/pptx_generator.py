from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from typing import List, Dict

# Try to import ImageGenerator, but don't fail if it's not available
try:
    from app.services.image_generator import ImageGenerator
    IMAGE_GENERATION_AVAILABLE = True
except ImportError:
    IMAGE_GENERATION_AVAILABLE = False
    ImageGenerator = None


class PPTXGenerator:
    def __init__(self, template: str = "modern", color_scheme: str = "blue", use_images: bool = False):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

        # Initialize image generator if requested AND available
        self.use_images = use_images and IMAGE_GENERATION_AVAILABLE
        if self.use_images:
            try:
                self.image_generator = ImageGenerator()
            except Exception as e:
                print(f"⚠️  Could not initialize image generator: {e}")
                self.image_generator = None
                self.use_images = False
        else:
            self.image_generator = None
        
        # Color schemes
        self.color_schemes = {
            "blue": {
                "primary": RGBColor(37, 99, 235),
                "secondary": RGBColor(59, 130, 246),
                "accent": RGBColor(96, 165, 250),
                "dark": RGBColor(30, 41, 59),
                "light": RGBColor(241, 245, 249),
                "text": RGBColor(15, 23, 42)
            },
            "green": {
                "primary": RGBColor(22, 163, 74),
                "secondary": RGBColor(34, 197, 94),
                "accent": RGBColor(74, 222, 128),
                "dark": RGBColor(20, 83, 45),
                "light": RGBColor(220, 252, 231),
                "text": RGBColor(1, 50, 32)
            },
            "purple": {
                "primary": RGBColor(147, 51, 234),
                "secondary": RGBColor(168, 85, 247),
                "accent": RGBColor(192, 132, 252),
                "dark": RGBColor(76, 29, 149),
                "light": RGBColor(243, 232, 255),
                "text": RGBColor(59, 7, 100)
            },
            "orange": {
                "primary": RGBColor(249, 115, 22),
                "secondary": RGBColor(251, 146, 60),
                "accent": RGBColor(253, 186, 116),
                "dark": RGBColor(124, 45, 18),
                "light": RGBColor(255, 237, 213),
                "text": RGBColor(67, 20, 7)
            },
            "dark": {
                "primary": RGBColor(51, 65, 85),
                "secondary": RGBColor(71, 85, 105),
                "accent": RGBColor(148, 163, 184),
                "dark": RGBColor(15, 23, 42),
                "light": RGBColor(248, 250, 252),
                "text": RGBColor(15, 23, 42)
            }
        }
        
        self.colors = self.color_schemes.get(color_scheme, self.color_schemes["blue"])
    
    def generate_presentation(self, slides_data: List[Dict], presentation_title: str, output_path: str):
        """Generate presentation"""
        
        print(f"\n🎨 Creating presentation...")
        
        # Title slide
        self.create_title_slide(presentation_title)
        
        # Content slides
        for idx, slide_data in enumerate(slides_data, 1):
            print(f"   📄 Creating slide {idx}/{len(slides_data)}: {slide_data.get('title', 'Untitled')}")
            
            is_chapter_divider = slide_data.get("is_chapter_divider", False)
            
            if is_chapter_divider:
                self.create_chapter_divider_slide(slide_data)
            else:
                # Use academic layout for all content slides
                self.create_academic_content_slide(slide_data)
        
        # Save presentation
        self.prs.save(output_path)
        print(f"   ✅ Saved: {output_path}\n")
    
    def create_title_slide(self, title: str):
        """Title slide with gradient"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Full-width gradient background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.colors["primary"]
        background.line.fill.background()
        
        # Accent bar
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.5), Inches(10), Inches(0.5)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.colors["accent"]
        accent.line.fill.background()
        
        # Main title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(2)
        )
        text_frame = title_box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(5), Inches(8), Inches(1)
        )
        text_frame = subtitle_box.text_frame
        
        p = text_frame.paragraphs[0]
        p.text = "Generated by EduSlide AI"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    def create_chapter_divider_slide(self, slide_data: Dict):
        """Chapter divider slide"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors["light"]
        bg.line.fill.background()
        
        # Colored accent
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(2), Inches(7.5)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.colors["primary"]
        accent.line.fill.background()
        
        # Title
        title_text = slide_data.get("title", "Chapter")
        title_text = title_text.replace('*', '').replace('**', '').replace('✓', '').replace('✔', '').strip()
        
        title_box = slide.shapes.add_textbox(
            Inches(2.5), Inches(3), Inches(6), Inches(1.5)
        )
        text_frame = title_box.text_frame
        
        p = text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.colors["dark"]
    
    def create_academic_content_slide(self, slide_data: Dict):
        """Academic layout: Title + Intro Paragraph + Bullet Points"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # White background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.line.fill.background()
        
        # TITLE - CLEANED, SIZE 20PT
        title_text = slide_data.get("title", "")
        
        # REMOVE ALL ASTERISKS AND SPECIAL CHARS
        title_text = title_text.replace('*', '').replace('**', '')
        title_text = title_text.replace('✓', '').replace('✔', '')
        title_text = title_text.strip()
        
        # Truncate title if too long
        if len(title_text) > 100:
            title_text = title_text[:97] + "..."
        
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(0.8)
        )
        text_frame = title_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        
        p = text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors["primary"]
        p.line_spacing = 1.1
        
        # CONTENT: Intro Paragraph + Bullets
        content_list = slide_data.get("content", [])
        
        if content_list and len(content_list) > 0:
            # FIRST ITEM = Introduction Paragraph
            intro_text = content_list[0]
            
            # Clean all special characters
            intro_text = intro_text.replace('*', '').replace('**', '')
            intro_text = intro_text.replace('✓', '').replace('✔', '')
            intro_text = intro_text.strip()
            
            # Truncate intro if too long
            if len(intro_text) > 280:
                intro_text = intro_text[:277] + "..."
            
            intro_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.4), Inches(9), Inches(1.2)
            )
            intro_frame = intro_box.text_frame
            intro_frame.word_wrap = True
            
            intro_p = intro_frame.paragraphs[0]
            intro_p.text = intro_text
            intro_p.font.size = Pt(14)
            intro_p.font.color.rgb = self.colors["text"]
            intro_p.line_spacing = 1.3
        
        # BULLET POINTS (NO CHECKMARKS)
        bullets = content_list[1:5] if len(content_list) > 1 else []
        
        if bullets:
            bullet_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.8), Inches(9), Inches(4.4)
            )
            bullet_frame = bullet_box.text_frame
            bullet_frame.word_wrap = True
            
            for i, bullet_text in enumerate(bullets):
                # Clean all special characters
                bullet_text = bullet_text.replace('*', '').replace('**', '')
                bullet_text = bullet_text.replace('✓', '').replace('✔', '')
                bullet_text = bullet_text.strip()
                
                # Truncate bullet if too long
                if len(bullet_text) > 140:
                    bullet_text = bullet_text[:137] + "..."
                
                if i == 0:
                    p = bullet_frame.paragraphs[0]
                else:
                    p = bullet_frame.add_paragraph()
                
                # SIMPLE BULLET - NO CHECKMARK
                p.text = f"• {bullet_text}"
                p.font.size = Pt(14)
                p.font.color.rgb = self.colors["text"]
                p.space_before = Pt(10)
                p.space_after = Pt(10)
                p.line_spacing = 1.2
                p.level = 0
        
        # Accent line
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(7.2),
            Inches(2), Inches(0.1)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = self.colors["primary"]
        accent_line.line.fill.background()